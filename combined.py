import os
import re
import sys
import io
import json
import base64
import datetime
import streamlit as st
import streamlit.components.v1 as components
from dotenv import load_dotenv

sys.path.insert(0, os.path.join(os.path.dirname(__file__), 'backend'))

from backend.auth import init_db, register_user, authenticate_user
from backend.ingest import ingest_text
from backend.rag import generate_answer

from youtube_transcript_api import YouTubeTranscriptApi
from youtube_transcript_api._errors import (
    NoTranscriptFound, TranscriptsDisabled,
    VideoUnavailable, CouldNotRetrieveTranscript, YouTubeRequestFailed,
)

from langchain_text_splitters import RecursiveCharacterTextSplitter
from langchain_community.vectorstores import Chroma
from langchain_classic.chains import RetrievalQA
from langchain_core.prompts import PromptTemplate
from langchain_core.documents import Document
from langchain_huggingface import HuggingFaceEmbeddings
from langchain_groq import ChatGroq

from pypdf import PdfReader
from pptx import Presentation

# ─────────────────────────────────────────────
# ENV & Constants
# ─────────────────────────────────────────────

load_dotenv()
GROQ_API_KEY = os.getenv("GROQ_API_KEY")
EMBED_MODEL  = "all-MiniLM-L6-v2"
DIGEST_STORE = os.path.join(os.path.dirname(os.path.abspath(__file__)), "digest_history.json")

# ─────────────────────────────────────────────
# Page config  (MUST come before any st.* call)
# ─────────────────────────────────────────────

st.set_page_config(
    page_title="Second Brain",
    page_icon="🧠",
    layout="wide",
    initial_sidebar_state="expanded",
)

# ─────────────────────────────────────────────
# Global CSS
# ─────────────────────────────────────────────

st.markdown("""
<link href="https://fonts.googleapis.com/css2?family=Playfair+Display:ital,wght@0,600;0,700;0,800;1,600&family=Outfit:wght@300;400;500;600;700&family=JetBrains+Mono:wght@400;500&display=swap" rel="stylesheet">

<style>
/* ════════════════════════════════════════════════════════
   DESIGN TOKENS — Obsidian × Warm Ivory × Amber
════════════════════════════════════════════════════════ */
:root {
  /* canvas */
  --bg:       #F9FAFB;
  --bg2:      #F3F4F6;
  --bg3:      #E5E7EB;
  --surface:  #ffffff;
  --glass:    rgba(255,255,255,0.8);

  /* ink */
  --ink:      #111827;
  --ink2:     #374151;
  --ink3:     #4B5563;
  --ink4:     #9CA3AF;

  /* sidebar */
  --sb:       #F9FAFB;
  --sb2:      #F3F4F6;
  --sb3:      #E5E7EB;
  --sb4:      #D1D5DB;
  --sb5:      #9CA3AF;
  --sbt:      #111827;
  --sbt2:     #374151;
  --sbt3:     #6B7280;

  /* accents (Blue primary) */
  --amber:    #2563EB;
  --amber-l:  #DBEAFE;
  --amber-m:  #60A5FA;
  --amber-d:  #1D4ED8;

  /* accents (Purple accent) */
  --teal:     #7C3AED;
  --teal-l:   #EDE9FE;
  --teal-m:   #A78BFA;

  --rose:     #EF4444;
  --rose-l:   #FEE2E2;
  --rose-m:   #F87171;

  --sky:      #0EA5E9;
  --sky-l:    #E0F2FE;
  --sky-m:    #7DD3FC;

  --plum:     #8B5CF6;
  --plum-l:   #EDE9FE;
  --plum-m:   #C4B5FD;

  /* radii */
  --r-sm:  8px;
  --r:     12px;
  --r-lg:  16px;
  --r-xl:  24px;

  /* shadows */
  --sh-xs:  0 1px 2px 0 rgba(0, 0, 0, 0.05);
  --sh-sm:  0 1px 3px 0 rgba(0, 0, 0, 0.1), 0 1px 2px -1px rgba(0, 0, 0, 0.1);
  --sh-md:  0 4px 6px -1px rgba(0, 0, 0, 0.1), 0 2px 4px -2px rgba(0, 0, 0, 0.1);
  --sh-lg:  0 10px 15px -3px rgba(0, 0, 0, 0.1), 0 4px 6px -4px rgba(0, 0, 0, 0.1);
}

/* ════════════════════════════════════════════════════════
   GLOBAL RESET & BASE
════════════════════════════════════════════════════════ */
*, *::before, *::after { box-sizing: border-box; }

html, body, [class*="css"] {
  font-family: 'Outfit', sans-serif !important;
  -webkit-font-smoothing: antialiased;
}

/* parchment canvas with subtle noise */
.stApp {
  background: var(--bg) !important;
}

#MainMenu, footer { visibility: hidden; }
header { visibility: hidden; }
header [data-testid="stSidebarCollapsedControl"] { visibility: visible !important; }

/* ════════════════════════════════════════════════════════
   SCROLLBAR
════════════════════════════════════════════════════════ */
::-webkit-scrollbar { width: 5px; height: 5px; }
::-webkit-scrollbar-track { background: transparent; }
::-webkit-scrollbar-thumb { background: var(--bg3); border-radius: 99px; }
::-webkit-scrollbar-thumb:hover { background: var(--ink4); }

/* ════════════════════════════════════════════════════════
   SIDEBAR — deep obsidian
════════════════════════════════════════════════════════ */
[data-testid="stSidebar"] {
  background: var(--sb) !important;
  border-right: 1px solid var(--sb3) !important;
  box-shadow: var(--sh-sm) !important;
}
[data-testid="stSidebar"] section[data-testid="stSidebarContent"] { padding: 0 !important; }

[data-testid="stSidebar"] *:not(button):not(input):not(textarea):not(select):not(option) {
  color: var(--sbt) !important;
  font-family: 'Outfit', sans-serif !important;
}

/* nav radio */
[data-testid="stSidebar"] .stRadio > div { gap: 4px !important; }
[data-testid="stSidebar"] .stRadio label {
  background: transparent !important;
  border-radius: var(--r-sm) !important;
  padding: 8px 12px !important;
  font-size: 14px !important;
  font-weight: 500 !important;
  cursor: pointer !important;
  transition: all 0.2s ease !important;
  color: var(--sbt2) !important;
  display: flex !important;
  align-items: center !important;
  gap: 8px !important;
}
[data-testid="stSidebar"] .stRadio label:hover {
  background: var(--sb3) !important;
  color: var(--sbt) !important;
}
[data-testid="stSidebar"] .stRadio [data-checked="true"] > label,
[data-testid="stSidebar"] .stRadio label[data-baseweb="radio"]:has(input:checked) {
  background: var(--amber-l) !important;
  color: var(--amber) !important;
  box-shadow: none !important;
  border-left: 4px solid var(--amber) !important;
}
[data-testid="stSidebar"] .stRadio [data-baseweb="radio"] > div:first-child { display: none !important; }

/* sidebar inputs */
[data-testid="stSidebar"] .stTextArea textarea,
[data-testid="stSidebar"] .stTextInput input {
  background: var(--surface) !important;
  border: 1px solid var(--sb4) !important;
  border-radius: var(--r-sm) !important;  
  color: var(--ink) !important;
  font-family: 'Outfit', sans-serif !important;
  font-size: 14px !important;
  transition: border-color 0.2s, box-shadow 0.2s !important;
}
[data-testid="stSidebar"] .stTextArea textarea:focus,
[data-testid="stSidebar"] .stTextInput input:focus {
  border-color: var(--amber) !important;
  box-shadow: 0 0 0 3px var(--amber-l) !important;
}
[data-testid="stSidebar"] .stTextArea textarea::placeholder,
[data-testid="stSidebar"] .stTextInput input::placeholder { color: var(--sbt3) !important; }

/* sidebar buttons */
[data-testid="stSidebar"] .stButton > button {
  font-family: 'Outfit', sans-serif !important;
  font-weight: 500 !important;
  font-size: 14px !important;
  border-radius: var(--r) !important;
  padding: 10px 16px !important;
  transition: all 0.2s ease !important;
  background: var(--surface) !important;
  border: 1px solid var(--sb4) !important;
  color: var(--sbt) !important;
}
[data-testid="stSidebar"] .stButton > button:hover {
  background: var(--bg2) !important;
  color: var(--teal) !important;
  border-color: var(--teal) !important;
}
[data-testid="stSidebar"] .stButton > button[kind="primary"] {
  background: var(--amber) !important;
  border-color: transparent !important;
  color: #fff !important;
  box-shadow: var(--sh-sm) !important;
}
[data-testid="stSidebar"] .stButton > button[kind="primary"]:hover {
  background: var(--teal) !important;
  transform: translateY(-1px) !important;
}

/* sidebar file uploader */
[data-testid="stSidebar"] [data-testid="stFileUploader"] {
  background: var(--sb2) !important;
  border: 1.5px dashed var(--sb5) !important;
  border-radius: var(--r) !important;
  transition: all 0.2s !important;
}
[data-testid="stSidebar"] [data-testid="stFileUploader"]:hover { border-color: var(--amber) !important; }
[data-testid="stSidebar"] [data-testid="stFileUploader"] label { display: none !important; }
[data-testid="stSidebar"] [data-testid="stFileUploader"] button {
  background: var(--sb3) !important;
  border: 1px solid var(--sb4) !important;
  color: var(--sbt) !important;
  border-radius: var(--r-sm) !important;
  font-family: 'Outfit', sans-serif !important;
  font-size: 12px !important;
}
[data-testid="stSidebar"] [data-testid="stFileUploaderDropzone"] { color: var(--sbt2) !important; }
[data-testid="stSidebar"] [data-testid="stFileUploaderDropzone"] small { color: var(--sbt3) !important; }

/* sidebar selectbox */
[data-testid="stSidebar"] [data-testid="stSelectbox"] > div > div {
  background: var(--sb2) !important;
  border: 1px solid var(--sb4) !important;
  border-radius: var(--r-sm) !important;
  color: var(--ink2) !important;
  font-family: 'Outfit', sans-serif !important;
  font-size: 13px !important;
}

[data-testid="stSidebar"] hr {
  border-color: var(--sb3) !important;
  border-width: 1px !important;
  margin: 12px 0 !important;
}

/* sidebar alerts */
[data-testid="stSidebar"] .stSuccess {
  background: rgba(13,158,130,0.12) !important;
  border-left: 3px solid var(--teal) !important;
  border-radius: 0 var(--r-sm) var(--r-sm) 0 !important;
  color: var(--teal) !important;
}
[data-testid="stSidebar"] .stError {
  background: rgba(214,63,86,0.12) !important;
  border-left: 3px solid var(--rose) !important;
  border-radius: 0 var(--r-sm) var(--r-sm) 0 !important;
  color: var(--rose) !important;
}
[data-testid="stSidebar"] .stWarning {
  background: rgba(212,134,10,0.12) !important;
  border-left: 3px solid var(--amber) !important;
  border-radius: 0 var(--r-sm) var(--r-sm) 0 !important;
  color: var(--amber) !important;
}
[data-testid="stSidebar"] .stInfo {
  background: rgba(114,69,232,0.12) !important;
  border-left: 3px solid var(--plum) !important;
  border-radius: 0 var(--r-sm) var(--r-sm) 0 !important;
  color: var(--plum) !important;
}
[data-testid="stSidebar"] .stSpinner > div { border-top-color: var(--amber) !important; }

/* ════════════════════════════════════════════════════════
   MAIN AREA — INPUTS
════════════════════════════════════════════════════════ */
.stTextArea textarea, .stTextInput input {
  background: var(--surface) !important;
  border: 1.5px solid var(--bg3) !important;
  border-radius: var(--r) !important;
  color: var(--ink) !important;
  font-family: 'Outfit', sans-serif !important;
  font-size: 14px !important;
  box-shadow: var(--sh-xs) !important;
  transition: border-color 0.15s, box-shadow 0.15s !important;
}
.stTextArea textarea:focus, .stTextInput input:focus {
  border-color: var(--amber) !important;
  box-shadow: 0 0 0 3px var(--amber-l), var(--sh-sm) !important;
}

/* ════════════════════════════════════════════════════════
   MAIN AREA — BUTTONS
════════════════════════════════════════════════════════ */
.main .stButton > button, section.main .stButton > button {
  font-family: 'Outfit', sans-serif !important;
  font-weight: 500 !important;
  font-size: 14px !important;
  border-radius: var(--r) !important;
  padding: 10px 20px !important;
  transition: all 0.2s ease !important;
  border: 1px solid var(--bg3) !important;
  background: var(--surface) !important;
  color: var(--ink2) !important;
  box-shadow: var(--sh-xs) !important;
}
.main .stButton > button:hover {
  background: var(--bg) !important;
  border-color: var(--teal) !important;
  color: var(--teal) !important;
  box-shadow: var(--sh-sm) !important;
}
.main .stButton > button:active { transform: translateY(0) !important; }
.main .stButton > button[kind="primary"] {
  background: var(--amber) !important;
  border-color: transparent !important;
  color: #fff !important;
  box-shadow: var(--sh-sm) !important;
}
.main .stButton > button[kind="primary"]:hover {
  background: var(--teal) !important;
  box-shadow: var(--sh-md) !important;
  transform: translateY(-1px) !important;
}

/* ════════════════════════════════════════════════════════
   CHAT MESSAGES
════════════════════════════════════════════════════════ */
[data-testid="stChatMessage"] {
  background: var(--surface) !important;
  border: 1px solid var(--bg3) !important;
  border-radius: var(--r-lg) !important;
  padding: 16px 20px !important;
  margin-bottom: 16px !important;
  font-size: 15px !important;
  box-shadow: var(--sh-sm) !important;
  transition: box-shadow 0.2s !important;
  animation: msg-in 0.25s ease !important;
}
@keyframes msg-in {
  from { opacity: 0; transform: translateY(6px); }
  to   { opacity: 1; transform: translateY(0); }
}
[data-testid="stChatMessage"]:hover { box-shadow: var(--sh-md) !important; }
[data-testid="stChatMessage"]:has([data-testid="chatAvatarIcon-user"]) {
  background: var(--bg2) !important;
  border-color: var(--bg3) !important;
}
[data-testid="stChatInput"] > div {
  background: var(--surface) !important;
  border: 1px solid var(--bg3) !important;
  border-radius: var(--r-lg) !important;
  box-shadow: var(--sh-md) !important;
  transition: border-color 0.2s, box-shadow 0.2s !important;
}
[data-testid="stChatInput"] > div:focus-within {
  border-color: var(--amber) !important;
  box-shadow: 0 0 0 3px var(--amber-l), var(--sh-md) !important;
}
[data-testid="stChatInput"] textarea {
  font-family: 'Outfit', sans-serif !important;
  font-size: 15px !important;
  color: var(--ink) !important;
}

/* ════════════════════════════════════════════════════════
   METRICS
════════════════════════════════════════════════════════ */
[data-testid="stMetric"] {
  background: var(--surface) !important;
  border: 1px solid var(--bg3) !important;
  border-radius: var(--r) !important;
  padding: 16px 20px !important;
  box-shadow: var(--sh-xs) !important;
  transition: box-shadow 0.2s, transform 0.2s !important;
}
[data-testid="stMetric"]:hover {
  box-shadow: var(--sh-sm) !important;
  transform: translateY(-1px) !important;
}
[data-testid="stMetricValue"] {
  font-family: 'Playfair Display', serif !important;
  font-size: 22px !important;
  font-weight: 700 !important;
  color: var(--amber) !important;
  letter-spacing: -0.5px !important;
}
[data-testid="stMetricLabel"] {
  font-size: 9.5px !important;
  font-weight: 600 !important;
  letter-spacing: 1.5px !important;
  text-transform: uppercase !important;
  color: var(--ink3) !important;
}

/* ════════════════════════════════════════════════════════
   EXPANDERS
════════════════════════════════════════════════════════ */
[data-testid="stExpander"] {
  background: var(--surface) !important;
  border: 1px solid var(--bg3) !important;
  border-radius: var(--r) !important;
  overflow: hidden !important;
  box-shadow: var(--sh-xs) !important;
  transition: box-shadow 0.2s !important;
}
[data-testid="stExpander"]:hover { box-shadow: var(--sh-sm) !important; }
[data-testid="stExpander"] summary {
  font-family: 'Outfit', sans-serif !important;
  font-weight: 600 !important;
  font-size: 13.5px !important;
  color: var(--ink) !important;
  padding: 14px 18px !important;
  transition: background 0.15s !important;
}
[data-testid="stExpander"] summary:hover { background: var(--bg2) !important; }

/* ════════════════════════════════════════════════════════
   ALERTS — main area
════════════════════════════════════════════════════════ */
.stSuccess {
  background: var(--teal-l) !important;
  border-left: 3px solid var(--teal) !important;
  border-radius: 0 var(--r) var(--r) 0 !important;
  color: var(--teal) !important;
  box-shadow: var(--sh-xs) !important;
}
.stError {
  background: var(--rose-l) !important;
  border-left: 3px solid var(--rose) !important;
  border-radius: 0 var(--r) var(--r) 0 !important;
  color: var(--rose) !important;
}
.stWarning {
  background: var(--amber-l) !important;
  border-left: 3px solid var(--amber) !important;
  border-radius: 0 var(--r) var(--r) 0 !important;
  color: var(--amber) !important;
}
.stInfo {
  background: var(--plum-l) !important;
  border-left: 3px solid var(--plum) !important;
  border-radius: 0 var(--r) var(--r) 0 !important;
  color: var(--plum) !important;
}

/* ════════════════════════════════════════════════════════
   DOWNLOAD BUTTON
════════════════════════════════════════════════════════ */
[data-testid="stDownloadButton"] > button {
  font-family: 'Outfit', sans-serif !important;
  font-weight: 600 !important;
  font-size: 13.5px !important;
  border-radius: var(--r-sm) !important;
  border: 1.5px solid var(--bg3) !important;
  background: var(--surface) !important;
  color: var(--ink2) !important;
  width: 100% !important;
  transition: all 0.18s !important;
  box-shadow: var(--sh-xs) !important;
}
[data-testid="stDownloadButton"] > button:hover {
  background: var(--bg2) !important;
  transform: translateY(-1px) !important;
  box-shadow: var(--sh-sm) !important;
}

/* ════════════════════════════════════════════════════════
   SELECTBOX (main)
════════════════════════════════════════════════════════ */
[data-testid="stSelectbox"] > div > div {
  background: var(--surface) !important;
  border: 1.5px solid var(--bg3) !important;
  border-radius: var(--r) !important;
  font-family: 'Outfit', sans-serif !important;
  font-size: 13.5px !important;
  color: var(--ink) !important;
  box-shadow: var(--sh-xs) !important;
}

/* ════════════════════════════════════════════════════════
   FILE UPLOADER (main)
════════════════════════════════════════════════════════ */
[data-testid="stFileUploader"] {
  background: var(--surface) !important;
  border: 2px dashed var(--bg3) !important;
  border-radius: var(--r-lg) !important;
  transition: all 0.2s !important;
}
[data-testid="stFileUploader"]:hover {
  border-color: var(--amber-m) !important;
  background: var(--amber-l) !important;
}

/* ════════════════════════════════════════════════════════
   SPINNER / DIVIDER
════════════════════════════════════════════════════════ */
.stSpinner > div { border-top-color: var(--amber) !important; }
hr {
  border: none !important;
  border-top: 1px solid var(--bg3) !important;
  margin: 20px 0 !important;
}

/* ════════════════════════════════════════════════════════
   SIDEBAR CUSTOM COMPONENTS
════════════════════════════════════════════════════════ */

/* Logo block */
.sb-logo {
  display: flex;
  align-items: center;
  gap: 13px;
  padding: 24px 20px 20px;
  border-bottom: 1px solid var(--sb3);
  margin-bottom: 4px;
}
.sb-blob {
  width: 40px; height: 40px;
  background: var(--amber);
  border-radius: 11px;
  display: flex; align-items: center; justify-content: center;
  font-size: 20px; flex-shrink: 0;
  box-shadow: 0 3px 12px rgba(212,134,10,0.4);
}
.sb-name {
  font-family: 'Playfair Display', serif;
  font-size: 19px; font-weight: 700;
  color: var(--sbt) !important;
  letter-spacing: -0.3px;
  line-height: 1.15;
}
.sb-tag {
  font-family: 'JetBrains Mono', monospace;
  font-size: 8.5px; letter-spacing: 2.5px;
  text-transform: uppercase;
  color: var(--sbt3) !important;
  margin-top: 2px;
}

/* Section label */
.sb-label {
  font-family: 'JetBrains Mono', monospace;
  font-size: 8.5px; letter-spacing: 2.5px;
  text-transform: uppercase;
  color: var(--sbt3) !important;
  padding: 16px 20px 6px;
  display: block;
}

/* Stat card */
.sb-stat {
  margin: 14px 14px 0;
  background: var(--sb2);
  border: 1px solid var(--sb3);
  border-radius: var(--r);
  padding: 16px 18px;
  position: relative;
  overflow: hidden;
}
.sb-stat::before {
  content: '';
  position: absolute; top: 0; left: 0; right: 0; height: 2px;
  background: linear-gradient(90deg, var(--amber), transparent);
}
.sb-stat-label {
  font-family: 'JetBrains Mono', monospace;
  font-size: 8.5px; letter-spacing: 1.8px;
  text-transform: uppercase;
  color: var(--sbt3) !important;
  margin-bottom: 7px;
}
.sb-stat-num {
  font-family: 'Playfair Display', serif;
  font-size: 32px; font-weight: 700;
  color: var(--amber-m) !important;
  letter-spacing: -1px; line-height: 1.1;
}
.sb-stat-sub { font-size: 11px; color: var(--sbt2) !important; margin-top: 4px; }
.sb-bar {
  height: 3px;
  background: var(--sb3);
  border-radius: 99px;
  margin-top: 12px;
  overflow: hidden;
}
.sb-fill {
  height: 100%;
  background: linear-gradient(90deg, var(--amber), var(--amber-m));
  border-radius: 99px;
  transition: width 0.6s cubic-bezier(.4,0,.2,1);
}
.sb-tags { display: flex; gap: 5px; flex-wrap: wrap; margin-top: 9px; }
.sb-t {
  font-family: 'JetBrains Mono', monospace;
  font-size: 8px; padding: 2px 8px;
  border-radius: 4px;
  border: 1px solid var(--sb3);
  color: var(--sbt2) !important;
  letter-spacing: 0.5px;
}

/* ════════════════════════════════════════════════════════
   PAGE HEADING
════════════════════════════════════════════════════════ */
.pg-head {
  display: flex;
  align-items: center;
  gap: 14px;
  padding-bottom: 20px;
  border-bottom: 1px solid var(--bg3);
  margin-bottom: 28px;
}
.pg-h1 {
  font-family: 'Playfair Display', serif;
  font-size: 30px; font-weight: 700;
  color: var(--ink);
  letter-spacing: -0.5px;
  margin: 0;
  line-height: 1.15;
}
.pg-pill {
  font-family: 'JetBrains Mono', monospace;
  font-size: 9px; font-weight: 700;
  letter-spacing: 1.5px;
  text-transform: uppercase;
  padding: 4px 12px;
  border-radius: 99px;
}
.pill-v  { background: var(--plum-l);  color: var(--plum);  border: 1px solid var(--plum-m); }
.pill-t  { background: var(--teal-l);  color: var(--teal);  border: 1px solid var(--teal-m); }
.pill-g  { background: var(--amber-l); color: var(--amber); border: 1px solid var(--amber-m); }
.pill-q  { background: var(--rose-l);  color: var(--rose);  border: 1px solid var(--rose-m); }
.pill-s  { background: var(--sky-l);   color: var(--sky);   border: 1px solid var(--sky-m); }
.pill-kb { background: var(--bg2);     color: var(--ink2);  border: 1px solid var(--bg3); }

/* ════════════════════════════════════════════════════════
   SUMMARY CARD
════════════════════════════════════════════════════════ */
.sum-card {
  background: var(--surface);
  border: 1px solid var(--bg3);
  border-radius: var(--r-lg);
  padding: 22px 24px;
  margin-bottom: 22px;
  box-shadow: var(--sh-sm);
  position: relative;
  overflow: hidden;
}
.sum-card::before {
  content: '';
  position: absolute; top: 0; left: 0; bottom: 0; width: 4px;
  background: linear-gradient(180deg, var(--amber), var(--amber-m));
}
.sum-head {
  display: flex; align-items: center;
  justify-content: space-between; margin-bottom: 12px;
}
.sum-title {
  font-family: 'Outfit', sans-serif;
  font-weight: 700; font-size: 14px; color: var(--ink);
}
.sum-chip {
  font-family: 'JetBrains Mono', monospace;
  font-size: 9px; font-weight: 700; letter-spacing: 1px;
  padding: 3px 10px; border-radius: 99px;
  background: var(--amber-l);
  border: 1px solid var(--amber-m);
  color: var(--amber);
}
.sum-body {
  font-size: 14px; color: var(--ink2);
  line-height: 1.8; font-style: italic;
  font-family: 'Outfit', sans-serif;
}

/* ════════════════════════════════════════════════════════
   SEPARATOR LABEL
════════════════════════════════════════════════════════ */
.sep-lbl {
  display: flex; align-items: center; gap: 12px;
  margin: 24px 0; color: var(--ink4);
  font-family: 'JetBrains Mono', monospace;
  font-size: 9px; letter-spacing: 2px; text-transform: uppercase;
}
.sep-lbl::before, .sep-lbl::after {
  content: ''; flex: 1; height: 1px; background: var(--bg3);
}

/* ════════════════════════════════════════════════════════
   FILE TYPE TAGS
════════════════════════════════════════════════════════ */
.ftag-row { display: flex; gap: 6px; margin-top: 8px; }
.ftag-pdf  {
  font-family: 'JetBrains Mono', monospace;
  font-size: 9px; font-weight: 700; letter-spacing: 1px;
  padding: 3px 10px; border-radius: 5px;
  background: var(--rose-l); color: var(--rose);
  border: 1px solid var(--rose-m);
}
.ftag-pptx {
  font-family: 'JetBrains Mono', monospace;
  font-size: 9px; font-weight: 700; letter-spacing: 1px;
  padding: 3px 10px; border-radius: 5px;
  background: var(--amber-l); color: var(--amber);
  border: 1px solid var(--amber-m);
}

/* ════════════════════════════════════════════════════════
   DIGEST
════════════════════════════════════════════════════════ */
.dg-head { display: flex; align-items: center; gap: 10px; margin: 24px 0 10px; }
.dg-icon {
  width: 28px; height: 28px; border-radius: var(--r-sm);
  display: flex; align-items: center; justify-content: center;
  font-size: 15px; flex-shrink: 0;
}
.dg-ht {
  font-family: 'Playfair Display', serif;
  font-size: 16px; font-weight: 700; letter-spacing: -0.2px;
}

/* ════════════════════════════════════════════════════════
   QUIZ COMPONENTS
════════════════════════════════════════════════════════ */
.quiz-question-card {
  background: var(--surface);
  border: 1px solid var(--rose-m);
  border-radius: var(--r-xl);
  padding: 28px 28px 24px;
  margin-bottom: 22px;
  position: relative;
  box-shadow: var(--sh-sm);
  overflow: hidden;
}
.quiz-question-card::after {
  content: '';
  position: absolute; bottom: 0; left: 0; right: 0; height: 3px;
  background: linear-gradient(90deg, var(--rose), var(--rose-m));
}
.quiz-q-label {
  font-family: 'JetBrains Mono', monospace;
  font-size: 9px; letter-spacing: 2px;
  text-transform: uppercase;
  color: var(--rose);
  margin-bottom: 14px;
  display: flex; align-items: center; gap: 8px;
}
.quiz-q-label::before {
  content: '';
  display: inline-block;
  width: 7px; height: 7px;
  border-radius: 50%;
  background: var(--rose);
}
.quiz-q-text {
  font-family: 'Playfair Display', serif;
  font-size: 20px; font-weight: 600;
  color: var(--ink); line-height: 1.5;
  letter-spacing: -0.2px;
}
.quiz-q-badge {
  position: absolute; top: 18px; right: 20px;
  font-family: 'JetBrains Mono', monospace;
  font-size: 9px; font-weight: 700; letter-spacing: 1px;
  padding: 3px 10px; border-radius: 99px;
  background: var(--rose-l);
  border: 1px solid var(--rose-m);
  color: var(--rose);
}

/* quiz result cards */
.quiz-result-correct {
  background: var(--teal-l);
  border: 1px solid var(--teal-m);
  border-radius: var(--r-lg);
  padding: 20px 22px; margin-top: 18px;
  box-shadow: var(--sh-xs);
}
.quiz-result-correct .quiz-result-icon { font-size: 24px; margin-bottom: 6px; }
.quiz-result-correct .quiz-result-title {
  font-family: 'Playfair Display', serif;
  font-weight: 700; font-size: 16px; color: var(--teal); margin-bottom: 8px;
}
.quiz-result-wrong {
  background: var(--rose-l);
  border: 1px solid var(--rose-m);
  border-radius: var(--r-lg);
  padding: 20px 22px; margin-top: 18px;
  box-shadow: var(--sh-xs);
}
.quiz-result-wrong .quiz-result-icon { font-size: 24px; margin-bottom: 6px; }
.quiz-result-wrong .quiz-result-title {
  font-family: 'Playfair Display', serif;
  font-weight: 700; font-size: 16px; color: var(--rose); margin-bottom: 8px;
}
.quiz-answer-box {
  background: var(--surface);
  border: 1px solid var(--bg3);
  border-radius: var(--r);
  padding: 14px 18px;
  font-size: 13.5px; color: var(--ink2);
  line-height: 1.75; margin-top: 10px;
}

/* score card */
.quiz-score-card {
  background: var(--bg2);
  border: 1px solid var(--bg3);
  border-radius: var(--r-xl);
  padding: 36px 32px 28px;
  text-align: center; margin-bottom: 28px;
  box-shadow: 0 8px 40px rgba(0,0,0,0.3);
  position: relative;
  overflow: hidden;
}
.quiz-score-card::before {
  content: '';
  position: absolute; top: 0; left: 0; right: 0; height: 2px;
  background: linear-gradient(90deg, var(--amber), var(--amber-m), var(--amber));
}
.quiz-score-big {
  font-family: 'Playfair Display', serif;
  font-size: 72px; font-weight: 800;
  color: var(--amber-m);
  letter-spacing: -3px; line-height: 1;
  text-shadow: 0 0 40px rgba(212,134,10,0.3);
}
.quiz-score-label {
  font-family: 'JetBrains Mono', monospace;
  font-size: 9px; letter-spacing: 3px;
  text-transform: uppercase;
  color: var(--ink4) !important;
  margin-top: 6px;
}
.quiz-score-sub { font-size: 14px; color: var(--ink3) !important; margin-top: 12px; }

/* progress bar */
.quiz-prog-wrap {
  background: var(--bg2);
  border-radius: 99px;
  height: 5px; overflow: hidden; margin-bottom: 24px;
}
.quiz-prog-fill {
  height: 100%; border-radius: 99px;
  background: linear-gradient(90deg, var(--rose), var(--rose-m));
  transition: width 0.5s cubic-bezier(.4,0,.2,1);
}

/* history item */
.qh-item {
  display: flex; align-items: flex-start; gap: 13px;
  padding: 14px 18px;
  background: var(--surface);
  border: 1px solid var(--bg3);
  border-radius: var(--r);
  margin-bottom: 8px;
  box-shadow: var(--sh-xs);
  transition: box-shadow 0.18s, transform 0.18s;
}
.qh-item:hover { box-shadow: var(--sh-sm); transform: translateX(2px); }
.qh-icon {
  width: 30px; height: 30px; border-radius: var(--r-sm);
  display: flex; align-items: center; justify-content: center;
  font-size: 15px; flex-shrink: 0; margin-top: 1px;
}
.qh-correct { background: var(--teal-l); }
.qh-wrong   { background: var(--rose-l); }
.qh-q  { font-size: 13px; font-weight: 600; color: var(--ink); margin-bottom: 3px; }
.qh-ua { font-size: 12px; color: var(--ink3); line-height: 1.5; }
.sq-stat-num {
  font-family: 'Playfair Display', serif;
  font-size: 26px; font-weight: 700;
  color: var(--rose-m) !important;
}

/* ════════════════════════════════════════════════════════
   SEMANTIC SEARCH
════════════════════════════════════════════════════════ */
.ss-live-badge {
  display: inline-flex; align-items: center; gap: 6px;
  font-family: 'JetBrains Mono', monospace;
  font-size: 9px; letter-spacing: 2px; text-transform: uppercase;
  color: var(--sky); margin-bottom: 12px;
}
.ss-live-dot {
  width: 6px; height: 6px; border-radius: 50%;
  background: var(--sky);
  animation: ss-pulse 1.4s ease-in-out infinite;
}
@keyframes ss-pulse {
  0%, 100% { opacity: 1; transform: scale(1); }
  50%       { opacity: 0.35; transform: scale(0.65); }
}
.ss-result-card {
  background: var(--surface);
  border: 1px solid var(--bg3);
  border-radius: var(--r-lg);
  padding: 18px 20px;
  margin-bottom: 12px;
  transition: all 0.2s cubic-bezier(.4,0,.2,1);
  position: relative; overflow: hidden;
  box-shadow: var(--sh-xs);
}
.ss-result-card:hover {
  border-color: var(--sky-m);
  box-shadow: var(--sh-md);
  transform: translateY(-2px);
}
.ss-result-card::before {
  content: '';
  position: absolute; left: 0; top: 0; bottom: 0; width: 4px;
  background: linear-gradient(180deg, var(--sky), var(--sky-m));
  border-radius: 4px 0 0 4px;
}
.ss-rank-badge {
  position: absolute; top: 14px; right: 14px;
  font-family: 'JetBrains Mono', monospace;
  font-size: 9px; letter-spacing: 1px;
  text-transform: uppercase;
  padding: 3px 9px; border-radius: 99px;
  background: var(--sky-l); color: var(--sky);
  border: 1px solid var(--sky-m);
}
.ss-score-bar-wrap {
  height: 3px; background: var(--bg2);
  border-radius: 99px; overflow: hidden; margin: 10px 0 12px;
}
.ss-score-bar {
  height: 100%; border-radius: 99px;
  background: linear-gradient(90deg, var(--sky), var(--sky-m));
  transition: width 0.4s ease;
}
.ss-snippet {
  font-size: 13.5px; color: var(--ink2);
  line-height: 1.8; font-family: 'Outfit', sans-serif;
}
.ss-snippet mark {
  background: var(--sky-l); color: var(--sky);
  border-radius: 3px; padding: 1px 3px; font-weight: 600;
}
.ss-meta { display: flex; align-items: center; gap: 8px; margin-top: 12px; flex-wrap: wrap; }
.ss-meta-tag {
  font-family: 'JetBrains Mono', monospace;
  font-size: 9px; letter-spacing: 1px; text-transform: uppercase;
  padding: 3px 9px; border-radius: 5px;
  background: var(--bg2); color: var(--ink3);
  border: 1px solid var(--bg3);
}
.ss-meta-score {
  font-family: 'JetBrains Mono', monospace;
  font-size: 9px; letter-spacing: 1px;
  color: var(--sky); margin-left: auto;
  font-weight: 600;
}
.ss-empty {
  text-align: center; padding: 48px 20px;
  color: var(--ink3); font-size: 14px;
}
.ss-empty-icon { font-size: 36px; margin-bottom: 12px; }
.ss-empty-msg  { font-family: 'Outfit', sans-serif; line-height: 1.7; }
.ss-stats-row { display: flex; gap: 10px; margin-bottom: 20px; flex-wrap: wrap; }
.ss-stat-chip {
  font-family: 'JetBrains Mono', monospace;
  font-size: 9.5px; letter-spacing: 0.8px;
  padding: 5px 13px; border-radius: 7px;
  background: var(--sky-l); color: var(--sky);
  border: 1px solid var(--sky-m);
}
.ss-sb-num {
  font-family: 'Playfair Display', serif;
  font-size: 26px; font-weight: 700;
  color: var(--sky-m) !important;
}

/* ════════════════════════════════════════════════════════
   KNOWLEDGE BASE — INGESTION LOG
════════════════════════════════════════════════════════ */
.il-metrics {
  display: grid; grid-template-columns: repeat(4, 1fr);
  gap: 12px; margin-bottom: 24px;
}
.il-metric {
  background: var(--surface);
  border: 1px solid var(--bg3);
  border-radius: var(--r);
  padding: 16px 18px; text-align: center;
  box-shadow: var(--sh-xs);
  transition: box-shadow 0.2s, transform 0.2s;
  position: relative; overflow: hidden;
}
.il-metric:hover { box-shadow: var(--sh-sm); transform: translateY(-2px); }
.il-metric::after {
  content: '';
  position: absolute; bottom: 0; left: 0; right: 0; height: 2px;
}
.il-m-pdf::after  { background: var(--rose); }
.il-m-pptx::after { background: var(--amber); }
.il-m-text::after { background: var(--plum); }
.il-m-url::after  { background: var(--teal); }
.il-metric-val {
  font-family: 'Playfair Display', serif;
  font-size: 28px; font-weight: 700;
  letter-spacing: -1px; line-height: 1.1;
}
.il-metric-lbl {
  font-family: 'JetBrains Mono', monospace;
  font-size: 8px; letter-spacing: 2px;
  text-transform: uppercase; color: var(--ink3); margin-top: 5px;
}
.il-m-pdf  .il-metric-val { color: var(--rose); }
.il-m-pptx .il-metric-val { color: var(--amber); }
.il-m-text .il-metric-val { color: var(--plum); }
.il-m-url  .il-metric-val { color: var(--teal); }

/* breakdown bar */
.il-breakdown {
  background: var(--surface);
  border: 1px solid var(--bg3);
  border-radius: var(--r);
  padding: 16px 18px; margin-bottom: 18px;
  box-shadow: var(--sh-xs);
}
.il-breakdown-title {
  font-family: 'JetBrains Mono', monospace;
  font-size: 9px; letter-spacing: 2px; text-transform: uppercase;
  color: var(--ink3); margin-bottom: 12px;
}
.il-bar-row { display: flex; align-items: center; gap: 12px; margin-bottom: 9px; }
.il-bar-label { font-size: 11.5px; font-weight: 600; color: var(--ink2); width: 38px; flex-shrink: 0; }
.il-bar-track {
  flex: 1; height: 7px; background: var(--bg2);
  border-radius: 99px; overflow: hidden;
}
.il-bar-fill { height: 100%; border-radius: 99px; transition: width 0.6s cubic-bezier(.4,0,.2,1); }
.il-bar-count {
  font-family: 'JetBrains Mono', monospace;
  font-size: 10px; color: var(--ink3);
  width: 24px; text-align: right;
}

/* item rows */
.il-item {
  display: flex; align-items: flex-start; gap: 13px;
  padding: 14px 16px;
  background: var(--surface);
  border: 1px solid var(--bg3);
  border-radius: var(--r);
  margin-bottom: 8px;
  transition: all 0.18s cubic-bezier(.4,0,.2,1);
  box-shadow: var(--sh-xs);
}
.il-item:hover {
  border-color: var(--amber-m);
  box-shadow: var(--sh-sm);
  transform: translateX(2px);
}
.il-type-icon {
  width: 36px; height: 36px;
  border-radius: var(--r-sm);
  display: flex; align-items: center; justify-content: center;
  font-size: 17px; flex-shrink: 0;
}
.il-icon-pdf  { background: var(--rose-l);  }
.il-icon-pptx { background: var(--amber-l); }
.il-icon-text { background: var(--plum-l);  }
.il-icon-url  { background: var(--teal-l);  }
.il-item-body { flex: 1; min-width: 0; }
.il-item-name {
  font-size: 13.5px; font-weight: 600; color: var(--ink);
  white-space: nowrap; overflow: hidden; text-overflow: ellipsis; margin-bottom: 3px;
}
.il-item-meta {
  font-family: 'JetBrains Mono', monospace;
  font-size: 9px; letter-spacing: 0.5px; color: var(--ink3);
}
.il-item-chunks {
  font-family: 'Playfair Display', serif;
  font-size: 16px; font-weight: 700;
  color: var(--amber); flex-shrink: 0; align-self: center;
}
.il-empty {
  text-align: center; padding: 40px 20px;
  color: var(--ink3); font-size: 14px;
  border: 2px dashed var(--bg3); border-radius: var(--r-lg);
}
.il-empty-icon { font-size: 32px; margin-bottom: 10px; }
.il-header { display: flex; align-items: center; justify-content: space-between; margin-bottom: 14px; }
.il-title {
  font-family: 'Playfair Display', serif;
  font-size: 18px; font-weight: 700; color: var(--ink); letter-spacing: -0.2px;
}
.il-total-chip {
  font-family: 'JetBrains Mono', monospace;
  font-size: 9px; letter-spacing: 1.5px; text-transform: uppercase;
  padding: 3px 11px; border-radius: 99px;
  background: var(--amber-l); border: 1px solid var(--amber-m); color: var(--amber);
}

/* ════════════════════════════════════════════════════════
   FILE VIEWER PANEL
════════════════════════════════════════════════════════ */
.fv-panel {
  background: var(--surface);
  border: 1px solid var(--amber-m);
  border-radius: var(--r-xl);
  padding: 26px 28px;
  margin: 4px 0 20px;
  box-shadow: var(--sh-md);
  animation: fv-drop-in 0.25s cubic-bezier(.4,0,.2,1);
  position: relative;
  overflow: hidden;
}
.fv-panel::before {
  content: '';
  position: absolute; top: 0; left: 0; right: 0; height: 3px;
  background: linear-gradient(90deg, var(--amber), var(--amber-m), transparent);
}
@keyframes fv-drop-in {
  from { opacity: 0; transform: translateY(-8px) scale(0.99); }
  to   { opacity: 1; transform: translateY(0) scale(1); }
}
.fv-header {
  display: flex; align-items: center;
  justify-content: space-between; margin-bottom: 18px;
  padding-bottom: 16px; border-bottom: 1px solid var(--bg3);
}
.fv-title {
  font-family: 'Playfair Display', serif;
  font-size: 17px; font-weight: 700;
  color: var(--ink);
  display: flex; align-items: center; gap: 10px;
}
.fv-badge {
  font-family: 'JetBrains Mono', monospace;
  font-size: 9px; font-weight: 700;
  letter-spacing: 1.5px; text-transform: uppercase;
  padding: 3px 10px; border-radius: 99px;
  background: var(--amber-l); border: 1px solid var(--amber-m); color: var(--amber);
}
.fv-meta-row { display: flex; gap: 8px; flex-wrap: wrap; margin-bottom: 18px; }
.fv-meta-chip {
  font-family: 'JetBrains Mono', monospace;
  font-size: 9px; letter-spacing: 1px; text-transform: uppercase;
  padding: 4px 11px; border-radius: var(--r-sm);
  background: var(--bg); border: 1px solid var(--bg3); color: var(--ink3);
}
.fv-section-head {
  font-family: 'JetBrains Mono', monospace;
  font-size: 9px; letter-spacing: 2.5px; text-transform: uppercase;
  color: var(--amber); margin: 20px 0 10px;
  display: flex; align-items: center; gap: 10px;
}
.fv-section-head::after {
  content: ''; flex: 1; height: 1px;
  background: linear-gradient(90deg, var(--amber-m), transparent);
}
.fv-summary-box {
  background: var(--bg);
  border: 1px solid var(--bg3);
  border-radius: var(--r);
  padding: 18px 20px;
  font-size: 14px; color: var(--ink2);
  line-height: 1.85; font-style: italic;
  margin-bottom: 4px;
  box-shadow: inset 0 1px 3px rgba(0,0,0,0.04);
}
.fv-pdf-wrap {
  background: var(--bg);
  border: 1px solid var(--bg3);
  border-radius: var(--r);
  overflow: hidden; margin-top: 10px;
  box-shadow: var(--sh-sm);
}
.fv-text-wrap {
  background: var(--bg);
  border: 1px solid var(--bg3);
  border-radius: var(--r);
  padding: 16px 18px;
  font-family: 'JetBrains Mono', monospace;
  font-size: 11.5px; color: var(--ink2);
  line-height: 1.75; max-height: 440px;
  overflow-y: auto; white-space: pre-wrap;
  margin-top: 10px;
  box-shadow: inset 0 1px 4px rgba(0,0,0,0.04);
}

/* ════════════════════════════════════════════════════════
   FOOTER
════════════════════════════════════════════════════════ */
.sb-footer {
  text-align: center;
  color: var(--ink3);
  font-family: 'JetBrains Mono', monospace;
  font-size: 9.5px; letter-spacing: 0.5px;
  padding: 24px 0 10px;
  border-top: 1px solid var(--bg3);
  margin-top: 40px;
}
.sb-footer b { color: var(--amber); font-weight: 700; }

</style>
""", unsafe_allow_html=True)


# ─────────────────────────────────────────────
# File Extraction
# ─────────────────────────────────────────────

def extract_text_from_pdf(file_bytes: bytes) -> str:
    reader = PdfReader(io.BytesIO(file_bytes))
    pages = []
    for i, page in enumerate(reader.pages):
        t = page.extract_text()
        if t and t.strip():
            pages.append(f"[Page {i+1}]\n{t.strip()}")
    return "\n\n".join(pages)


def extract_text_from_pptx(file_bytes: bytes) -> str:
    prs = Presentation(io.BytesIO(file_bytes))
    slides = []
    for i, slide in enumerate(prs.slides):
        parts = []
        for shape in slide.shapes:
            if shape.has_text_frame:
                for para in shape.text_frame.paragraphs:
                    line = " ".join(r.text for r in para.runs).strip()
                    if line:
                        parts.append(line)
        if slide.has_notes_slide:
            notes = slide.notes_slide.notes_text_frame.text.strip()
            if notes:
                parts.append(f"[Notes]: {notes}")
        if parts:
            slides.append(f"[Slide {i+1}]\n" + "\n".join(parts))
    return "\n\n".join(slides)


def extract_text(file_bytes: bytes, fext: str) -> str:
    return extract_text_from_pdf(file_bytes) if fext == "pdf" else extract_text_from_pptx(file_bytes)


# ─────────────────────────────────────────────
# LLM helpers
# ─────────────────────────────────────────────

def get_llm() -> ChatGroq:
    return ChatGroq(groq_api_key=GROQ_API_KEY, model_name="openai/gpt-oss-120b", temperature=0.2)


@st.cache_resource(show_spinner="Loading embedding model...")
def get_embeddings() -> HuggingFaceEmbeddings:
    return HuggingFaceEmbeddings(model_name=EMBED_MODEL, encode_kwargs={"normalize_embeddings": True})


@st.cache_resource(show_spinner="Loading sentence transformer...")
def get_st_embedder():
    from sentence_transformers import SentenceTransformer
    return SentenceTransformer("all-MiniLM-L6-v2")


def summarise_with_llm(llm: ChatGroq, text: str, doc_type: str = "document") -> str:
    prompt = PromptTemplate(
        input_variables=["doc_type", "text"],
        template="You are a helpful assistant. Summarise the following {doc_type} clearly and concisely. Include all key points, main topics, and important details.\n\n{doc_type} content:\n{text}\n\nSummary:",
    )
    truncated = text[:6000] + ("\n\n[... content truncated ...]" if len(text) > 6000 else "")
    return llm.invoke(prompt.format(doc_type=doc_type, text=truncated)).content


# ─────────────────────────────────────────────
# Knowledge Digest Engine
# ─────────────────────────────────────────────

PROBE_QUERIES = [
    "main concepts and key ideas", "important facts and insights",
    "definitions and explanations", "conclusions and takeaways",
    "processes and how things work", "examples and case studies",
    "people organizations and entities mentioned",
    "problems challenges and solutions",
    "data statistics and numbers", "recommendations and action items",
]

DIGEST_PROMPT_TEMPLATE = """You are an expert knowledge curator and analyst.

Below is a representative sample of content from a personal knowledge base.

Generate a comprehensive **{period} Knowledge Digest** in structured Markdown with these sections:

## 📌 Executive Summary
## 🧠 Key Concepts & Ideas
## 🔗 Relationships & Connections
## 💡 Notable Insights
## 📚 Topics Covered
## ❓ Questions to Explore
## 🎯 Action Items & Recommendations

---

Retrieved knowledge base content ({chunk_count} chunks sampled):

{context}

---

Generate the {period} Knowledge Digest now. Be thorough and specific. Reference actual content — do not hallucinate.
"""


def fetch_semantic_sample(chroma_dir: str, k_per_query: int = 8) -> tuple[list[str], int]:
    import chromadb
    try:
        backend_dir      = os.path.join(os.path.dirname(os.path.abspath(__file__)), "backend")
        resolved_db_path = os.path.normpath(os.path.join(backend_dir, "..", "chroma_db"))
        chroma_host      = os.getenv("CHROMA_HOST")
        client           = chromadb.HttpClient() if chroma_host else chromadb.PersistentClient(path=resolved_db_path)
        try:
            col = client.get_collection(name="notes")
        except Exception:
            raise ValueError(f'ChromaDB collection "notes" not found. Ingest at least one document first.')
        if col.count() == 0:
            raise ValueError("The knowledge base is empty. Ingest some documents first.")
        embedder  = get_st_embedder()
        seen, chunks = set(), []
        k = min(k_per_query, max(1, col.count()))
        for query in PROBE_QUERIES:
            results = col.query(query_embeddings=[embedder.encode(query).tolist()], n_results=k, include=["documents"])
            for doc_text in results["documents"][0]:
                key = hash(doc_text[:200])
                if key not in seen:
                    seen.add(key)
                    chunks.append(doc_text)
        return chunks, len(chunks)
    except ValueError:
        raise
    except Exception as e:
        raise RuntimeError(f"Could not read from ChromaDB: {e}")


def generate_knowledge_digest(period: str, chroma_dir: str) -> dict:
    llm             = get_llm()
    chunks, count   = fetch_semantic_sample(chroma_dir, k_per_query=8)
    if not count:
        raise ValueError("No content found. Ingest some documents first.")
    parts, total = [], 0
    for i, chunk in enumerate(chunks):
        entry = f"--- Chunk {i+1} ---\n{chunk}"
        if total + len(entry) > 12000:
            break
        parts.append(entry)
        total += len(entry)
    prompt  = DIGEST_PROMPT_TEMPLATE.format(period=period, chunk_count=len(parts), context="\n\n".join(parts))
    content = llm.invoke(prompt).content
    return {"period": period, "generated_at": datetime.datetime.now().strftime("%Y-%m-%d %H:%M"), "chunk_count": len(parts), "content": content}


def load_digest_history() -> list[dict]:
    if os.path.exists(DIGEST_STORE):
        try:
            with open(DIGEST_STORE) as f:
                return json.load(f)
        except Exception:
            return []
    return []


def save_digest(digest: dict) -> None:
    history = load_digest_history()
    history.insert(0, digest)
    history = history[:20]
    try:
        with open(DIGEST_STORE, "w") as f:
            json.dump(history, f, indent=2)
    except Exception:
        pass


# ─────────────────────────────────────────────
# YouTube Helpers
# ─────────────────────────────────────────────

def extract_video_id(url: str):
    match = re.search(r"(?:v=|youtu\.be/|shorts/|embed/)([a-zA-Z0-9_-]{11})", url)
    return match.group(1) if match else None


def fetch_transcript(video_id: str) -> tuple[str | None, str | None]:
    try:
        ytt_api         = YouTubeTranscriptApi()
        transcript_list = ytt_api.list(video_id)
        available       = list(transcript_list)
        if not available:
            return None, "No transcripts of any kind are available for this video."
        transcript_obj = None
        warning        = None
        try:
            transcript_obj = transcript_list.find_manually_created_transcript(["en"])
        except NoTranscriptFound:
            pass
        if transcript_obj is None:
            try:
                transcript_obj = transcript_list.find_generated_transcript(["en"])
            except NoTranscriptFound:
                pass
        if transcript_obj is None:
            for t in available:
                if not t.is_generated and t.is_translatable:
                    try:
                        transcript_obj = t.translate("en")
                        warning = f"No English transcript found. Auto-translated from **{t.language}** ({t.language_code}) to English."
                        break
                    except Exception:
                        continue
        if transcript_obj is None:
            for t in available:
                if t.is_generated and t.is_translatable:
                    try:
                        transcript_obj = t.translate("en")
                        warning = f"No English transcript found. Auto-translated from **{t.language}** ({t.language_code}) to English."
                        break
                    except Exception:
                        continue
        if transcript_obj is None:
            transcript_obj = available[0]
            warning = f"No English transcript available. Using **{transcript_obj.language}** ({transcript_obj.language_code}) as-is. Q&A accuracy may be lower."
        data = transcript_obj.fetch()
        return " ".join(t.text for t in data), warning
    except TranscriptsDisabled:
        return None, "Transcripts are disabled by the video owner."
    except VideoUnavailable:
        return None, "This video is unavailable or private."
    except YouTubeRequestFailed as e:
        return None, f"YouTube blocked the request (403 Forbidden). Details: {e}"
    except CouldNotRetrieveTranscript as e:
        return None, f"Could not retrieve transcript: {e}"
    except Exception as e:
        return None, f"Unexpected error fetching transcript: {e}"


def build_youtube_qa_chain(transcript_text: str):
    docs   = [Document(page_content=transcript_text)]
    chunks = RecursiveCharacterTextSplitter(chunk_size=1000, chunk_overlap=200).split_documents(docs)
    vs     = Chroma.from_documents(chunks, get_embeddings())
    llm    = get_llm()
    prompt = PromptTemplate(
        input_variables=["context", "question"],
        template="Answer using ONLY the context below.\n\nContext:\n{context}\n\nQuestion:\n{question}\n\nIf not found, say: \"Not found in transcript.\"",
    )
    chain = RetrievalQA.from_chain_type(llm=llm, retriever=vs.as_retriever(search_kwargs={"k": 5}), chain_type_kwargs={"prompt": prompt})
    return chain, llm


# ─────────────────────────────────────────────
# Quiz Engine
# ─────────────────────────────────────────────

QUIZ_GEN_PROMPT = """You are an expert quiz generator. Based on the knowledge base content below, generate exactly {n} quiz questions.

Rules:
- Questions must be answerable strictly from the provided content
- Mix question types: factual recall, conceptual understanding, cause-and-effect
- Each question must have one clear correct answer
- Do NOT include multiple-choice options — open-ended questions only
- Keep questions concise (one sentence each)

Return ONLY valid JSON (no markdown, no backticks), exactly like this:
[
  {{
    "question": "What is ...?",
    "answer": "The correct answer is ..."
  }},
  ...
]

Knowledge base content:
{context}

Return the JSON array now:"""


QUIZ_EVAL_PROMPT = """You are a strict but fair quiz evaluator.

Question: {question}
Correct Answer (from knowledge base): {correct_answer}
User's Answer: {user_answer}

Evaluate whether the user's answer is correct. The answer doesn't need to be word-for-word identical — it just needs to capture the key idea correctly.

Respond in ONLY valid JSON (no markdown, no backticks):
{{
  "is_correct": true or false,
  "feedback": "One sentence explaining why correct or wrong.",
  "correct_answer": "The full correct answer from the knowledge base."
}}"""


def generate_quiz_questions(n: int = 5) -> list[dict]:
    chunks, count = fetch_semantic_sample("", k_per_query=6)
    if not count:
        raise ValueError("Knowledge base is empty. Ingest some documents first.")
    parts, total = [], 0
    for chunk in chunks:
        if total + len(chunk) > 8000:
            break
        parts.append(chunk)
        total += len(chunk)
    context = "\n\n---\n\n".join(parts)
    llm     = get_llm()
    prompt  = QUIZ_GEN_PROMPT.format(n=n, context=context)
    raw     = llm.invoke(prompt).content.strip()
    raw = re.sub(r"^```(?:json)?", "", raw, flags=re.MULTILINE).strip()
    raw = re.sub(r"```$", "", raw, flags=re.MULTILINE).strip()
    questions = json.loads(raw)
    return questions[:n]


def evaluate_answer(question: str, correct_answer: str, user_answer: str) -> dict:
    llm    = get_llm()
    prompt = QUIZ_EVAL_PROMPT.format(
        question=question,
        correct_answer=correct_answer,
        user_answer=user_answer,
    )
    raw = llm.invoke(prompt).content.strip()
    raw = re.sub(r"^```(?:json)?", "", raw, flags=re.MULTILINE).strip()
    raw = re.sub(r"```$", "", raw, flags=re.MULTILINE).strip()
    return json.loads(raw)


# ─────────────────────────────────────────────
# Semantic Search Engine
# ─────────────────────────────────────────────

def semantic_search(query: str, k: int = 8) -> list[dict]:
    import chromadb
    backend_dir      = os.path.join(os.path.dirname(os.path.abspath(__file__)), "backend")
    resolved_db_path = os.path.normpath(os.path.join(backend_dir, "..", "chroma_db"))
    chroma_host      = os.getenv("CHROMA_HOST")
    client           = chromadb.HttpClient() if chroma_host else chromadb.PersistentClient(path=resolved_db_path)
    try:
        col = client.get_collection(name="notes")
    except Exception:
        raise ValueError('ChromaDB collection "notes" not found. Ingest at least one document first.')
    if col.count() == 0:
        raise ValueError("Knowledge base is empty. Ingest some documents first.")
    embedder  = get_st_embedder()
    query_vec = embedder.encode(query).tolist()
    k_actual  = min(k, col.count())
    results = col.query(
        query_embeddings=[query_vec],
        n_results=k_actual,
        include=["documents", "distances", "metadatas"],
    )
    docs      = results["documents"][0]
    distances = results["distances"][0]
    metadatas = results["metadatas"][0]
    hits = []
    for doc, dist, meta in zip(docs, distances, metadatas):
        similarity = max(0.0, 1.0 - dist / 2.0)
        hits.append({"text": doc, "score": round(similarity * 100, 1), "metadata": meta or {}})
    hits.sort(key=lambda x: x["score"], reverse=True)
    return hits


def highlight_query(text: str, query: str, max_len: int = 280) -> str:
    lower_text  = text.lower()
    lower_query = query.lower()
    terms       = [t for t in lower_query.split() if len(t) > 2]
    start = 0
    for term in terms:
        pos = lower_text.find(term)
        if pos != -1:
            start = max(0, pos - 80)
            break
    snippet = text[start : start + max_len]
    if start > 0:
        snippet = "…" + snippet
    if start + max_len < len(text):
        snippet = snippet + "…"
    import html as html_mod
    snippet = html_mod.escape(snippet)
    for term in sorted(terms, key=len, reverse=True):
        pattern = re.compile(re.escape(html_mod.escape(term)), re.IGNORECASE)
        snippet = pattern.sub(lambda m: f"<mark>{m.group(0)}</mark>", snippet)
    return snippet


# ─────────────────────────────────────────────
# Session State Init
# ─────────────────────────────────────────────

defaults = {
    "logged_in": False, "username": None,
    "messages": [], "ingested_chunks": 0, "active_tab": "docs",
    "yt_processed": False, "yt_summary": "", "yt_chat_history": [],
    "yt_qa_chain": None, "yt_transcript_warn": None,
    "file_summary": "", "file_name": "", "file_ingested": False,
    "digest_result": None, "digest_history": [],
    "quiz_questions": [], "quiz_index": 0, "quiz_user_answer": "",
    "quiz_submitted": False, "quiz_eval": None, "quiz_history": [],
    "quiz_active": False, "quiz_finished": False, "quiz_score": 0,
    "ss_query": "", "ss_results": [], "ss_k": 6, "ss_searched": False,
    "ingest_log": [],       # list of {name, type, chunks, chars, timestamp, summary, file_bytes, raw_text}
    "kb_view_item": None,   # index of currently open file viewer panel
}
for k, v in defaults.items():
    if k not in st.session_state:
        st.session_state[k] = v

if not st.session_state.digest_history:
    st.session_state.digest_history = load_digest_history()

init_db()

if not st.session_state.logged_in:
    st.markdown("""
    <div style="display: flex; flex-direction: column; align-items: center; justify-content: center; margin-top: 10vh;">
      <div class="sb-blob" style="width: 60px; height: 60px; font-size: 30px; margin-bottom: 20px;">🧠</div>
      <div class="pg-h1" style="margin-bottom: 30px;">Second Brain Login</div>
    </div>
    """, unsafe_allow_html=True)

    _, col_auth, _ = st.columns([1, 1, 1])
    with col_auth:
        tab_login, tab_reg = st.tabs(["🔐 Login", "📝 Register"])
        
        with tab_login:
            login_user = st.text_input("Username", key="login_user")
            login_pass = st.text_input("Password", type="password", key="login_pass")
            if st.button("Login", type="primary", use_container_width=True):
                if authenticate_user(login_user, login_pass):
                    st.session_state.logged_in = True
                    st.session_state.username = login_user
                    st.rerun()
                else:
                    st.error("Invalid username or password.")
                    
        with tab_reg:
            reg_user = st.text_input("Choose Username", key="reg_user")
            reg_pass = st.text_input("Choose Password", type="password", key="reg_pass")
            if st.button("Register", use_container_width=True):
                success, msg = register_user(reg_user, reg_pass)
                if success:
                    st.success(msg + " You can now switch to Login.")
                else:
                    st.error(msg)
                    
    st.stop()


# ─────────────────────────────────────────────
# SIDEBAR
# ─────────────────────────────────────────────

with st.sidebar:

    st.markdown("""
    <div class="sb-logo">
      <div class="sb-blob">🧠</div>
      <div>
        <div class="sb-name">Second Brain</div>
        <div class="sb-tag">Knowledge OS</div>
      </div>
    </div>
    """, unsafe_allow_html=True)

    st.markdown('<span class="sb-label">Workspace</span>', unsafe_allow_html=True)
    mode = st.radio(
        "nav",
        ["📄  Document RAG", "🎥  YouTube RAG", "📅  Knowledge Digest", "🧩  Quiz Mode", "🔍  Semantic Search", "📦  Knowledge Base"],
        index={"docs": 0, "youtube": 1, "digest": 2, "quiz": 3, "search": 4, "kb": 5}.get(st.session_state.active_tab, 0),
        label_visibility="collapsed",
    )
    if "Document" in mode:    st.session_state.active_tab = "docs"
    elif "YouTube" in mode:   st.session_state.active_tab = "youtube"
    elif "Digest"  in mode:   st.session_state.active_tab = "digest"
    elif "Quiz"    in mode:   st.session_state.active_tab = "quiz"
    elif "Semantic" in mode:  st.session_state.active_tab = "search"
    else:                     st.session_state.active_tab = "kb"

    st.divider()

    # ── Docs sidebar ──────────────────────────────────────────────────────────
    if st.session_state.active_tab == "docs":
        st.markdown('<span class="sb-label">Input method</span>', unsafe_allow_html=True)
        doc_mode = st.radio("input", ["✏️  Paste Text", "📎  Upload File"],
                            horizontal=True, label_visibility="collapsed")

        if "Paste" in doc_mode:
            input_text = st.text_area("Content", height=150,
                                      placeholder="Paste articles, notes, research…",
                                      label_visibility="collapsed")
            if st.button("📤  Ingest Text", use_container_width=True, type="primary"):
                if input_text.strip():
                    with st.spinner("Processing and storing..."):
                        try:
                            r = ingest_text(input_text)
                            chunks = r.get("chunks", 0)
                            st.session_state.ingested_chunks += chunks
                            st.session_state.ingest_log.insert(0, {
                                "name":       f"Text snippet ({len(input_text)} chars)",
                                "type":       "text",
                                "chunks":     chunks,
                                "chars":      len(input_text),
                                "timestamp":  datetime.datetime.now().strftime("%d %b %Y, %H:%M"),
                                "summary":    "",
                                "file_bytes": None,
                                "raw_text":   input_text,
                            })
                            st.success(f"✅  {chunks} chunks stored.")
                        except Exception as e:
                            st.error(f"❌  {e}")
                else:
                    st.warning("Please enter some text first.")
        else:
            st.markdown("""
            <div class="ftag-row">
              <span class="ftag-pdf">PDF</span>
              <span class="ftag-pptx">PPTX</span>
            </div>
            """, unsafe_allow_html=True)
            uploaded_file = st.file_uploader("Choose a file:", type=["pdf", "pptx"], label_visibility="collapsed")
            if uploaded_file:
                fext       = uploaded_file.name.rsplit(".", 1)[-1].lower()
                doc_type   = "PDF document" if fext == "pdf" else "PowerPoint presentation"
                file_bytes = uploaded_file.read()
                c1, c2 = st.columns(2)
                with c1:
                    if st.button("🤖  Summarise", use_container_width=True, type="primary"):
                        with st.spinner("Extracting & summarising..."):
                            try:
                                raw = extract_text(file_bytes, fext)
                                if not raw.strip():
                                    st.error("No readable text found.")
                                else:
                                    st.session_state.file_summary  = summarise_with_llm(get_llm(), raw, doc_type)
                                    st.session_state.file_name     = uploaded_file.name
                                    st.session_state.file_ingested = False
                                    st.success("✅  Summary ready!")
                            except Exception as e:
                                st.error(f"❌  {e}")
                with c2:
                    if st.button("📥  Ingest", use_container_width=True):
                        with st.spinner("Ingesting..."):
                            try:
                                raw = extract_text(file_bytes, fext)
                                if not raw.strip():
                                    st.error("No readable text found.")
                                else:
                                    r = ingest_text(raw)
                                    chunks = r.get("chunks", 0)
                                    st.session_state.ingested_chunks += chunks
                                    st.session_state.file_ingested = True
                                    # Auto-generate summary for KB viewer
                                    try:
                                        auto_summary = summarise_with_llm(get_llm(), raw, doc_type)
                                    except Exception:
                                        auto_summary = ""
                                    st.session_state.ingest_log.insert(0, {
                                        "name":       uploaded_file.name,
                                        "type":       fext,
                                        "chunks":     chunks,
                                        "chars":      len(raw),
                                        "timestamp":  datetime.datetime.now().strftime("%d %b %Y, %H:%M"),
                                        "summary":    auto_summary,
                                        "file_bytes": file_bytes if fext == "pdf" else None,
                                        "raw_text":   raw,
                                    })
                                    st.success(f"✅  {chunks} chunks!")
                            except Exception as e:
                                st.error(f"❌  {e}")

        st.divider()
        chunks = st.session_state.ingested_chunks
        st.markdown(f"""
        <div class="sb-stat">
          <div class="sb-stat-label">Knowledge base</div>
          <div class="sb-stat-num">{chunks}</div>
          <div class="sb-stat-sub">chunks stored this session</div>
          <div class="sb-bar"><div class="sb-fill" style="width:{min(chunks, 200) / 2}%"></div></div>
          <div class="sb-tags"><span class="sb-t">session</span><span class="sb-t">chromadb</span></div>
        </div>
        """, unsafe_allow_html=True)
        st.markdown("")
        if st.button("🗑️  Clear Chat", use_container_width=True):
            st.session_state.messages = []
            st.rerun()

    # ── YouTube sidebar ───────────────────────────────────────────────────────
    elif st.session_state.active_tab == "youtube":
        st.markdown('<span class="sb-label">Video URL</span>', unsafe_allow_html=True)
        video_url = st.text_input("url", placeholder="https://www.youtube.com/watch?v=...", label_visibility="collapsed")
        if st.button("⚙️  Process Video", use_container_width=True, type="primary"):
            if not video_url.strip():
                st.warning("Please enter a URL.")
            else:
                vid = extract_video_id(video_url)
                if not vid:
                    st.error("Invalid YouTube URL.")
                else:
                    with st.spinner("Fetching transcript..."):
                        transcript, tw = fetch_transcript(vid)
                    if not transcript:
                        st.error(f"❌  {tw}")
                    else:
                        if tw:
                            st.warning(f"⚠️  {tw}")
                        with st.spinner("Building index & summarising..."):
                            chain, llm = build_youtube_qa_chain(transcript)
                            p = PromptTemplate(input_variables=["text"],
                                              template="Summarise this YouTube transcript clearly:\n\n{text}\n\nSummary:")
                            summary = llm.invoke(p.format(text=transcript[:5000])).content
                        st.session_state.update({
                            "yt_qa_chain": chain, "yt_summary": summary,
                            "yt_transcript_warn": tw, "yt_processed": True,
                            "yt_chat_history": [],
                        })
                        st.session_state.ingest_log.insert(0, {
                            "name":       video_url.strip(),
                            "type":       "url",
                            "chunks":     len(transcript.split()) // 200 + 1,
                            "chars":      len(transcript),
                            "timestamp":  datetime.datetime.now().strftime("%d %b %Y, %H:%M"),
                            "summary":    summary,
                            "file_bytes": None,
                            "raw_text":   transcript,
                        })
                        st.success("✅  Video processed!")
        if st.session_state.yt_processed:
            st.divider()
            if st.button("🗑️  Clear YouTube Chat", use_container_width=True):
                st.session_state.yt_chat_history = []
                st.rerun()

    # ── Digest sidebar ────────────────────────────────────────────────────────
    elif st.session_state.active_tab == "digest":
        st.markdown('<span class="sb-label">Digest settings</span>', unsafe_allow_html=True)
        digest_period = st.selectbox("Period", ["Weekly", "Monthly", "Custom Range"], label_visibility="collapsed")
        backend_dir      = os.path.join(os.path.dirname(os.path.abspath(__file__)), "backend")
        resolved_db_path = os.path.normpath(os.path.join(backend_dir, "..", "chroma_db"))
        st.caption(f"Knowledge base: `{resolved_db_path}`")
        chunks = st.session_state.ingested_chunks
        st.markdown(f"""
        <div class="sb-stat">
          <div class="sb-stat-label">Session chunks</div>
          <div class="sb-stat-num">{chunks}</div>
          <div class="sb-bar"><div class="sb-fill" style="width:{min(chunks,200)/2}%"></div></div>
        </div>
        """, unsafe_allow_html=True)
        st.divider()
        if st.button(f"⚡  Generate {digest_period} Digest", use_container_width=True, type="primary"):
            with st.spinner("Running semantic search across knowledge base..."):
                try:
                    d = generate_knowledge_digest(digest_period, resolved_db_path)
                    st.session_state.digest_result = d
                    save_digest(d)
                    st.session_state.digest_history = load_digest_history()
                    st.success("✅  Digest generated!")
                except ValueError as e:
                    st.error(f"⚠️  {e}")
                except Exception as e:
                    st.error(f"❌  Error: {e}")
        st.caption(f"💾  {len(st.session_state.digest_history)} past digest(s) saved")

    # ── Quiz sidebar ──────────────────────────────────────────────────────────
    elif st.session_state.active_tab == "quiz":
        st.markdown('<span class="sb-label">Quiz settings</span>', unsafe_allow_html=True)
        quiz_n = st.selectbox("Number of questions", [3, 5, 7, 10], index=1, label_visibility="collapsed")
        total_q   = len(st.session_state.quiz_questions)
        correct_q = st.session_state.quiz_score
        st.markdown(f"""
        <div class="sb-stat">
          <div class="sb-stat-label">Quiz score</div>
          <div class="sq-stat-num">{correct_q} / {total_q if total_q else "—"}</div>
          <div class="sb-stat-sub">correct answers this session</div>
          <div class="sb-bar">
            <div class="sb-fill" style="width:{int(correct_q/total_q*100) if total_q else 0}%; background: var(--bg2);"></div>
          </div>
        </div>
        """, unsafe_allow_html=True)
        st.divider()
        if st.button("🚀  Start New Quiz", use_container_width=True, type="primary"):
            with st.spinner("Generating questions from your knowledge base..."):
                try:
                    qs = generate_quiz_questions(n=quiz_n)
                    st.session_state.quiz_questions  = qs
                    st.session_state.quiz_index      = 0
                    st.session_state.quiz_submitted  = False
                    st.session_state.quiz_eval       = None
                    st.session_state.quiz_history    = []
                    st.session_state.quiz_active     = True
                    st.session_state.quiz_finished   = False
                    st.session_state.quiz_score      = 0
                    st.rerun()
                except ValueError as e:
                    st.error(f"⚠️  {e}")
                except Exception as e:
                    st.error(f"❌  Error: {e}")
        if st.session_state.quiz_active and not st.session_state.quiz_finished:
            if st.button("⏹  End Quiz", use_container_width=True):
                st.session_state.quiz_finished = True
                st.rerun()
        if st.session_state.quiz_history:
            st.divider()
            if st.button("🗑️  Clear Quiz History", use_container_width=True):
                for k in ["quiz_questions","quiz_index","quiz_submitted","quiz_eval",
                          "quiz_history","quiz_active","quiz_finished","quiz_score"]:
                    st.session_state[k] = defaults[k]
                st.rerun()

    # ── Semantic Search sidebar ───────────────────────────────────────────────
    else:
        st.markdown('<span class="sb-label">Search settings</span>', unsafe_allow_html=True)
        ss_k = st.selectbox("Results to show", [3, 5, 8, 10, 15], index=1, label_visibility="collapsed")
        st.session_state.ss_k = ss_k
        n_results = len(st.session_state.ss_results)
        st.markdown(f"""
        <div class="sb-stat">
          <div class="sb-stat-label">Last search</div>
          <div class="ss-sb-num">{n_results}</div>
          <div class="sb-stat-sub">chunks matched</div>
          <div class="sb-bar">
            <div class="sb-fill" style="width:{min(n_results * 10, 100)}%; background: var(--bg2);"></div>
          </div>
        </div>
        """, unsafe_allow_html=True)
        st.divider()
        if st.session_state.ss_searched and st.button("🗑️  Clear Results", use_container_width=True):
            st.session_state.ss_results  = []
            st.session_state.ss_query    = ""
            st.session_state.ss_searched = False
            st.rerun()


# ─────────────────────────────────────────────
# MAIN CONTENT
# ─────────────────────────────────────────────

col_welc1, col_welc2 = st.columns([4, 1])
with col_welc1:
    st.markdown(f'<div class="pg-h1" style="font-size: 20px; color: var(--amber); margin-bottom: 20px; border-bottom: 1px solid var(--bg3); padding-bottom: 10px;">👋 Welcome, {st.session_state.username}</div>', unsafe_allow_html=True)
with col_welc2:
    if st.button("🚪 Logout", use_container_width=True):
        st.session_state.logged_in = False
        st.session_state.username = None
        st.rerun()

# ══════════════════════════════════════════════
# TAB: Document RAG
# ══════════════════════════════════════════════
if st.session_state.active_tab == "docs":

    st.markdown("""
    <div class="pg-head">
      <div class="pg-h1">Document Q&amp;A</div>
      <div class="pg-pill pill-v">RAG</div>
    </div>
    """, unsafe_allow_html=True)

    if st.session_state.file_summary:
        fext_d = st.session_state.file_name.rsplit(".", 1)[-1].upper() if st.session_state.file_name else ""
        icon   = "📄" if fext_d == "PDF" else "📊"
        st.markdown(f"""
        <div class="sum-card">
          <div class="sum-head">
            <div class="sum-title">{icon} Summary — {st.session_state.file_name}</div>
            <span class="sum-chip">AI Summary</span>
          </div>
          <div class="sum-body">{st.session_state.file_summary}</div>
        </div>
        """, unsafe_allow_html=True)
        if not st.session_state.file_ingested:
            st.warning("⚠️  Not ingested yet — click **Ingest** in the sidebar to enable Q&A over this file.")
        else:
            st.success("✅  Ingested and ready for Q&A.")
        if st.button("✖  Dismiss Summary"):
            st.session_state.file_summary  = ""
            st.session_state.file_name     = ""
            st.session_state.file_ingested = False
            st.rerun()

    st.markdown('<div class="sep-lbl">conversation</div>', unsafe_allow_html=True)

    for msg in st.session_state.messages:
        with st.chat_message(msg["role"]):
            st.markdown(msg["content"])

    if prompt := st.chat_input("Ask anything about your documents…"):
        st.session_state.messages.append({"role": "user", "content": prompt})
        with st.chat_message("user"):
            st.markdown(prompt)
        with st.chat_message("assistant"):
            with st.spinner("Thinking…"):
                try:
                    resp = generate_answer(prompt)
                    st.markdown(resp)
                    st.session_state.messages.append({"role": "assistant", "content": resp})
                except Exception as e:
                    err = f"❌  Error: {e}"
                    st.error(err)
                    st.session_state.messages.append({"role": "assistant", "content": err})


# ══════════════════════════════════════════════
# TAB: YouTube RAG
# ══════════════════════════════════════════════
elif st.session_state.active_tab == "youtube":

    st.markdown("""
    <div class="pg-head">
      <div class="pg-h1">YouTube Q&amp;A</div>
      <div class="pg-pill pill-t">RAG</div>
    </div>
    """, unsafe_allow_html=True)

    if not st.session_state.yt_processed:
        st.info("👈  Paste a YouTube URL in the sidebar and hit **Process Video** to begin.")
    else:
        if st.session_state.yt_transcript_warn:
            st.warning(f"⚠️  {st.session_state.yt_transcript_warn}")
        st.markdown(f"""
        <div class="sum-card">
          <div class="sum-head">
            <div class="sum-title">📌 Video Summary</div>
            <span class="sum-chip">AI Summary</span>
          </div>
          <div class="sum-body">{st.session_state.yt_summary}</div>
        </div>
        """, unsafe_allow_html=True)
        st.markdown('<div class="sep-lbl">ask about the video</div>', unsafe_allow_html=True)
        for role, msg in st.session_state.yt_chat_history:
            with st.chat_message("user" if role == "You" else "assistant"):
                st.markdown(msg)
        if q := st.chat_input("Ask something about the video…"):
            st.session_state.yt_chat_history.append(("You", q))
            with st.chat_message("user"):
                st.markdown(q)
            with st.chat_message("assistant"):
                with st.spinner("Thinking…"):
                    try:
                        ans = st.session_state.yt_qa_chain.invoke({"query": q})["result"]
                        st.markdown(ans)
                        st.session_state.yt_chat_history.append(("AI", ans))
                    except Exception as e:
                        err = f"❌  Error: {e}"
                        st.error(err)
                        st.session_state.yt_chat_history.append(("AI", err))


# ══════════════════════════════════════════════
# TAB: Knowledge Digest
# ══════════════════════════════════════════════
elif st.session_state.active_tab == "digest":

    st.markdown("""
    <div class="pg-head">
      <div class="pg-h1">Knowledge Digest</div>
      <div class="pg-pill pill-g">Digest</div>
    </div>
    """, unsafe_allow_html=True)

    with st.expander("ℹ️  How this works", expanded=False):
        st.markdown("""
**The digest engine works in three steps:**

1. **Semantic Sampling** — 10 broad probe queries pull up to ~80 unique chunks from your ChromaDB vector store via similarity search.
2. **Context Assembly** — deduplicates and assembles chunks into a ~12 000-char context window.
3. **LLM Synthesis** — Groq identifies key concepts, cross-document relationships, notable insights, open questions, and action items.
        """)

    st.divider()

    if st.session_state.digest_result:
        d = st.session_state.digest_result
        c1, c2, c3 = st.columns(3)
        c1.metric("Period",    d["period"])
        c2.metric("Generated", d["generated_at"])
        c3.metric("Chunks",    d["chunk_count"])
        st.divider()
        st.markdown(d["content"])
        st.divider()
        st.download_button(
            "⬇️  Download Digest (.md)",
            data=d["content"],
            file_name=f"digest_{d['period'].lower()}_{d['generated_at'].replace(' ','_').replace(':','-')}.md",
            mime="text/markdown",
            use_container_width=True,
        )
    else:
        st.info("👈  Select a period and click **Generate Digest** in the sidebar.\n\nMake sure you've ingested at least a few documents first.")

    if st.session_state.digest_history:
        st.divider()
        st.markdown("""
        <div class="dg-head">
          <div class="dg-icon" style="background:var(--gl);">🗂️</div>
          <div class="dg-ht" style="color:var(--g);">Past Digests</div>
        </div>
        """, unsafe_allow_html=True)
        st.caption("The last 20 generated digests are saved locally.")
        for i, past in enumerate(st.session_state.digest_history):
            label = f"**{past['period']} Digest** — {past['generated_at']}  ·  {past['chunk_count']} chunks"
            with st.expander(label, expanded=False):
                st.markdown(past["content"])
                st.download_button(
                    "⬇️  Download",
                    data=past["content"],
                    file_name=f"digest_{past['period'].lower()}_{past['generated_at'].replace(' ','_').replace(':','-')}.md",
                    mime="text/markdown",
                    key=f"dl_past_{i}",
                )


# ══════════════════════════════════════════════
# TAB: Quiz Mode
# ══════════════════════════════════════════════
elif st.session_state.active_tab == "quiz":

    st.markdown("""
    <div class="pg-head">
      <div class="pg-h1">Quiz Mode</div>
      <div class="pg-pill pill-q">Test yourself</div>
    </div>
    """, unsafe_allow_html=True)

    if not st.session_state.quiz_active and not st.session_state.quiz_finished:
        st.info("👈  Choose the number of questions in the sidebar and click **Start New Quiz** to begin.\n\nQuiz questions are generated directly from your ingested knowledge base.")

    elif st.session_state.quiz_finished:
        total   = len(st.session_state.quiz_history)
        correct = st.session_state.quiz_score
        pct     = int(correct / total * 100) if total else 0
        if pct == 100:
            grade, msg = "🏆", "Perfect score! Outstanding knowledge!"
        elif pct >= 80:
            grade, msg = "🥇", "Excellent work! You know this material well."
        elif pct >= 60:
            grade, msg = "🥈", "Good effort! Review the missed questions below."
        elif pct >= 40:
            grade, msg = "🥉", "Keep studying — you're getting there!"
        else:
            grade, msg = "📚", "Needs more review. Re-read your ingested material."
        st.markdown(f"""
        <div class="quiz-score-card">
          <div style="font-size:36px; margin-bottom:8px;">{grade}</div>
          <div class="quiz-score-big">{pct}%</div>
          <div class="quiz-score-label">Final Score</div>
          <div class="quiz-score-sub">{correct} out of {total} correct &nbsp;·&nbsp; {msg}</div>
        </div>
        """, unsafe_allow_html=True)
        st.markdown(f"""
        <div class="quiz-prog-wrap">
          <div class="quiz-prog-fill" style="width:{pct}%"></div>
        </div>
        """, unsafe_allow_html=True)
        st.markdown('<div class="sep-lbl">question review</div>', unsafe_allow_html=True)
        for i, h in enumerate(st.session_state.quiz_history):
            icon  = "✅" if h["is_correct"] else "❌"
            cls   = "qh-correct" if h["is_correct"] else "qh-wrong"
            st.markdown(f"""
            <div class="qh-item">
              <div class="qh-icon {cls}">{icon}</div>
              <div style="flex:1; min-width:0;">
                <div class="qh-q">Q{i+1}. {h['question']}</div>
                <div class="qh-ua"><b>Your answer:</b> {h['user_answer'] or '<em>No answer given</em>'}</div>
                {"" if h['is_correct'] else f'<div class="qh-ua" style="color:#c44020;"><b>Correct answer:</b> {h["correct_answer"]}</div>'}
                <div class="qh-ua" style="color:var(--i3); margin-top:2px; font-style:italic;">{h['feedback']}</div>
              </div>
            </div>
            """, unsafe_allow_html=True)
        st.divider()
        if st.button("🔄  Retake / New Quiz", type="primary"):
            for k in ["quiz_questions","quiz_index","quiz_submitted","quiz_eval",
                      "quiz_history","quiz_active","quiz_finished","quiz_score"]:
                st.session_state[k] = defaults[k]
            st.rerun()

    elif st.session_state.quiz_active:
        questions = st.session_state.quiz_questions
        idx       = st.session_state.quiz_index
        total_q   = len(questions)
        if idx >= total_q:
            st.session_state.quiz_finished = True
            st.session_state.quiz_active   = False
            st.rerun()
        current   = questions[idx]
        answered  = len(st.session_state.quiz_history)
        correct   = st.session_state.quiz_score
        pct_done  = int(answered / total_q * 100)
        m1, m2, m3 = st.columns(3)
        m1.metric("Question",  f"{idx + 1} / {total_q}")
        m2.metric("Correct",   correct)
        m3.metric("Remaining", total_q - answered)
        st.markdown(f"""
        <div class="quiz-prog-wrap">
          <div class="quiz-prog-fill" style="width:{pct_done}%"></div>
        </div>
        """, unsafe_allow_html=True)
        st.markdown(f"""
        <div class="quiz-question-card">
          <div class="quiz-q-label">Question {idx + 1}</div>
          <div class="quiz-q-text">{current['question']}</div>
          <div class="quiz-q-badge">🧩 Quiz</div>
        </div>
        """, unsafe_allow_html=True)
        if not st.session_state.quiz_submitted:
            user_ans = st.text_area("Your answer", placeholder="Type your answer here…", height=100,
                                    key=f"quiz_ans_{idx}", label_visibility="collapsed")
            col_sub, col_skip = st.columns([3, 1])
            with col_sub:
                if st.button("✅  Submit Answer", type="primary", use_container_width=True):
                    if not user_ans.strip():
                        st.warning("Please type an answer before submitting.")
                    else:
                        with st.spinner("Evaluating your answer…"):
                            try:
                                result = evaluate_answer(
                                    question=current["question"],
                                    correct_answer=current["answer"],
                                    user_answer=user_ans,
                                )
                                st.session_state.quiz_eval      = result
                                st.session_state.quiz_submitted = True
                                if result.get("is_correct"):
                                    st.session_state.quiz_score += 1
                                st.session_state.quiz_history.append({
                                    "question":       current["question"],
                                    "user_answer":    user_ans,
                                    "correct_answer": result.get("correct_answer", current["answer"]),
                                    "is_correct":     result.get("is_correct", False),
                                    "feedback":       result.get("feedback", ""),
                                })
                                st.rerun()
                            except Exception as e:
                                st.error(f"❌  Evaluation error: {e}")
            with col_skip:
                if st.button("⏭  Skip", use_container_width=True):
                    st.session_state.quiz_history.append({
                        "question": current["question"], "user_answer": "",
                        "correct_answer": current["answer"], "is_correct": False, "feedback": "Skipped.",
                    })
                    st.session_state.quiz_index    += 1
                    st.session_state.quiz_submitted = False
                    st.session_state.quiz_eval      = None
                    if st.session_state.quiz_index >= total_q:
                        st.session_state.quiz_finished = True
                        st.session_state.quiz_active   = False
                    st.rerun()
        else:
            result = st.session_state.quiz_eval
            if result and result.get("is_correct"):
                st.markdown(f"""
                <div class="quiz-result-correct">
                  <div class="quiz-result-icon">🎉</div>
                  <div class="quiz-result-title">Correct!</div>
                  <div class="quiz-answer-box">{result.get('feedback', '')}</div>
                </div>
                """, unsafe_allow_html=True)
            else:
                st.markdown(f"""
                <div class="quiz-result-wrong">
                  <div class="quiz-result-icon">💡</div>
                  <div class="quiz-result-title">Not quite right</div>
                  <div class="quiz-answer-box">
                    <strong>Feedback:</strong> {result.get('feedback', '') if result else ''}<br><br>
                    <strong>Correct answer from your knowledge base:</strong><br>
                    {result.get('correct_answer', current['answer']) if result else current['answer']}
                  </div>
                </div>
                """, unsafe_allow_html=True)
            is_last   = (idx + 1 >= total_q)
            btn_label = "🏁  Finish Quiz" if is_last else "➡️  Next Question"
            if st.button(btn_label, type="primary", use_container_width=True):
                st.session_state.quiz_index    += 1
                st.session_state.quiz_submitted = False
                st.session_state.quiz_eval      = None
                if is_last:
                    st.session_state.quiz_finished = True
                    st.session_state.quiz_active   = False
                st.rerun()


# ══════════════════════════════════════════════
# TAB: Semantic Search
# ══════════════════════════════════════════════
elif st.session_state.active_tab == "search":

    st.markdown("""
    <div class="pg-head">
      <div class="pg-h1">Semantic Search</div>
      <div class="pg-pill pill-s">Vector</div>
    </div>
    """, unsafe_allow_html=True)

    col_inp, col_btn = st.columns([5, 1])
    with col_inp:
        query_input = st.text_input(
            "search", value=st.session_state.ss_query,
            placeholder="Type a concept, keyword, or phrase…",
            label_visibility="collapsed", key="ss_input_box",
        )
    with col_btn:
        search_btn = st.button("🔍  Search", type="primary", use_container_width=True)

    do_search = search_btn or (
        query_input.strip() and query_input.strip() != st.session_state.ss_query
    )
    if do_search and query_input.strip():
        with st.spinner("Searching your knowledge base…"):
            try:
                hits = semantic_search(query_input.strip(), k=st.session_state.ss_k)
                st.session_state.ss_results  = hits
                st.session_state.ss_query    = query_input.strip()
                st.session_state.ss_searched = True
                st.rerun()
            except ValueError as e:
                st.error(f"⚠️  {e}")
            except Exception as e:
                st.error(f"❌  Search error: {e}")

    if not st.session_state.ss_searched:
        st.markdown("""
        <div class="ss-empty">
          <div class="ss-empty-icon">🔍</div>
          <div class="ss-empty-msg">
            Type any word or concept above and press <strong>Enter</strong> or click <strong>Search</strong>.<br>
            The most semantically similar chunks from your knowledge base will appear instantly.
          </div>
        </div>
        """, unsafe_allow_html=True)
    elif not st.session_state.ss_results:
        st.markdown("""
        <div class="ss-empty">
          <div class="ss-empty-icon">🕳️</div>
          <div class="ss-empty-msg">No results found. Try a different search term or ingest more documents.</div>
        </div>
        """, unsafe_allow_html=True)
    else:
        hits  = st.session_state.ss_results
        query = st.session_state.ss_query
        top   = hits[0]["score"] if hits else 0
        avg_score = round(sum(h["score"] for h in hits) / len(hits), 1) if hits else 0
        st.markdown(f"""
        <div class="ss-stats-row">
          <span class="ss-stat-chip">🔍 Query: <strong>{query}</strong></span>
          <span class="ss-stat-chip">📦 {len(hits)} results</span>
          <span class="ss-stat-chip">⭐ Top match: {top}%</span>
          <span class="ss-stat-chip">📊 Avg: {avg_score}%</span>
        </div>
        """, unsafe_allow_html=True)
        st.markdown("""
        <div class="ss-live-badge">
          <div class="ss-live-dot"></div>
          Ranked by semantic similarity
        </div>
        """, unsafe_allow_html=True)
        for i, hit in enumerate(hits):
            score   = hit["score"]
            text    = hit["text"]
            meta    = hit["metadata"]
            snippet = highlight_query(text, query, max_len=300)
            bar_w   = int(score)
            source  = meta.get("source", meta.get("filename", meta.get("file", "")))
            src_tag = f'<span class="ss-meta-tag">📄 {source}</span>' if source else ""
            words   = len(text.split())
            wc_tag  = f'<span class="ss-meta-tag">~{words} words</span>'
            st.markdown(f"""
            <div class="ss-result-card">
              <div class="ss-rank-badge">#{i+1}</div>
              <div class="ss-score-bar-wrap"><div class="ss-score-bar" style="width:{bar_w}%"></div></div>
              <div class="ss-snippet">{snippet}</div>
              <div class="ss-meta">{src_tag}{wc_tag}<span class="ss-meta-tag">chunk {i+1}</span><span class="ss-meta-score">{score}% match</span></div>
            </div>
            """, unsafe_allow_html=True)
            with st.expander(f"View full chunk #{i+1}", expanded=False):
                st.markdown(f"**Similarity score:** `{score}%`")
                st.markdown(text)


# ══════════════════════════════════════════════
# TAB: Knowledge Base
# ══════════════════════════════════════════════
elif st.session_state.active_tab == "kb":

    log = st.session_state.ingest_log
    n_pdf  = sum(1 for x in log if x["type"] == "pdf")
    n_pptx = sum(1 for x in log if x["type"] == "pptx")
    n_text = sum(1 for x in log if x["type"] == "text")
    n_url  = sum(1 for x in log if x["type"] == "url")
    total_chunks = sum(x["chunks"] for x in log)
    total_chars  = sum(x["chars"]  for x in log)

    st.markdown("""
    <div class="pg-head">
      <div class="pg-h1">Knowledge Base</div>
      <div class="pg-pill pill-v">Ingestion Log</div>
    </div>
    """, unsafe_allow_html=True)

    st.markdown(f"""
    <div class="il-metrics">
      <div class="il-metric il-m-pdf">
        <div class="il-metric-val">{n_pdf}</div>
        <div class="il-metric-lbl">PDFs</div>
      </div>
      <div class="il-metric il-m-pptx">
        <div class="il-metric-val">{n_pptx}</div>
        <div class="il-metric-lbl">Slides</div>
      </div>
      <div class="il-metric il-m-text">
        <div class="il-metric-val">{n_text}</div>
        <div class="il-metric-lbl">Texts</div>
      </div>
      <div class="il-metric il-m-url">
        <div class="il-metric-val">{n_url}</div>
        <div class="il-metric-lbl">URLs</div>
      </div>
    </div>
    """, unsafe_allow_html=True)

    c1, c2, c3 = st.columns(3)
    c1.metric("Total Items",  len(log))
    c2.metric("Total Chunks", total_chunks)
    c3.metric("Total Chars",  f"{total_chars:,}")

    if log:
        st.divider()

        max_n = max(n_pdf, n_pptx, n_text, n_url, 1)
        def pct(n): return int(n / max_n * 100)

        st.markdown(f"""
        <div class="il-breakdown">
          <div class="il-breakdown-title">Breakdown by type</div>
          <div class="il-bar-row">
            <div class="il-bar-label">PDF</div>
            <div class="il-bar-track"><div class="il-bar-fill" style="width:{pct(n_pdf)}%; background:#e05533;"></div></div>
            <div class="il-bar-count">{n_pdf}</div>
          </div>
          <div class="il-bar-row">
            <div class="il-bar-label">PPTX</div>
            <div class="il-bar-track"><div class="il-bar-fill" style="width:{pct(n_pptx)}%; background:var(--g);"></div></div>
            <div class="il-bar-count">{n_pptx}</div>
          </div>
          <div class="il-bar-row">
            <div class="il-bar-label">Text</div>
            <div class="il-bar-track"><div class="il-bar-fill" style="width:{pct(n_text)}%; background:var(--v);"></div></div>
            <div class="il-bar-count">{n_text}</div>
          </div>
          <div class="il-bar-row">
            <div class="il-bar-label">URL</div>
            <div class="il-bar-track"><div class="il-bar-fill" style="width:{pct(n_url)}%; background:var(--t);"></div></div>
            <div class="il-bar-count">{n_url}</div>
          </div>
        </div>
        """, unsafe_allow_html=True)

        st.markdown("""
        <div class="il-header" style="margin-top:4px;">
          <div class="il-title">All Ingested Items</div>
        </div>
        """, unsafe_allow_html=True)

        type_meta = {
            "pdf":  ("📄", "il-icon-pdf",  "PDF"),
            "pptx": ("📊", "il-icon-pptx", "PPTX"),
            "text": ("✏️", "il-icon-text", "Text"),
            "url":  ("🎥", "il-icon-url",  "YouTube"),
        }

        for idx_item, item in enumerate(log):
            t    = item["type"]
            icon, icon_cls, label = type_meta.get(t, ("📎", "il-icon-text", t.upper()))
            name = item["name"]
            if t == "url":
                vid_match    = re.search(r"(?:v=|youtu\.be/)([a-zA-Z0-9_-]{11})", name)
                display_name = f"youtube.com/watch?v={vid_match.group(1)}" if vid_match else name
            else:
                display_name = name
            chars_k    = f"{item['chars'] // 1000}k" if item['chars'] >= 1000 else str(item['chars'])
            is_viewable = t in ("pdf", "pptx", "text")
            is_open     = st.session_state.kb_view_item == idx_item

            # Row: item card + view button
            col_item, col_view = st.columns([6, 1])
            with col_item:
                st.markdown(f"""
                <div class="il-item" style="{'border-color:var(--vm); background:var(--vl);' if is_open else ''}">
                  <div class="il-type-icon {icon_cls}">{icon}</div>
                  <div class="il-item-body">
                    <div class="il-item-name" title="{name}">{display_name}</div>
                    <div class="il-item-meta">{label} &nbsp;·&nbsp; {chars_k} chars &nbsp;·&nbsp; {item['timestamp']}</div>
                  </div>
                  <div class="il-item-chunks">{item['chunks']}<span style="font-size:9px; font-weight:400; color:var(--i3); margin-left:2px;">chunks</span></div>
                </div>
                """, unsafe_allow_html=True)
            with col_view:
                if is_viewable:
                    btn_label = "✖ Close" if is_open else "📂 View"
                    btn_type  = "primary" if is_open else "secondary"
                    if st.button(btn_label, key=f"view_btn_{idx_item}", use_container_width=True,
                                 type=btn_type if btn_type != "secondary" else "secondary"):
                        st.session_state.kb_view_item = None if is_open else idx_item
                        st.rerun()

            # ── File Viewer Panel ──────────────────────────────────────────────
            if is_open:
                fext_v   = item["type"]
                summary  = item.get("summary", "")
                raw_text = item.get("raw_text", "")
                fb       = item.get("file_bytes")
                file_ext_upper = fext_v.upper()

                st.markdown(f"""
                <div class="fv-panel">
                  <div class="fv-header">
                    <div class="fv-title">
                      {icon} {display_name}
                    </div>
                    <span class="fv-badge">{file_ext_upper}</span>
                  </div>
                  <div class="fv-meta-row">
                    <span class="fv-meta-chip">📦 {item['chunks']} chunks</span>
                    <span class="fv-meta-chip">📝 {chars_k} chars</span>
                    <span class="fv-meta-chip">🕒 {item['timestamp']}</span>
                  </div>
                """, unsafe_allow_html=True)

                # ── Summary section ──
                st.markdown('<div class="fv-section-head">AI Summary</div>', unsafe_allow_html=True)

                if summary:
                    st.markdown(f'<div class="fv-summary-box">{summary}</div>', unsafe_allow_html=True)
                else:
                    st.info("No summary generated yet for this item.")
                    if raw_text:
                        if st.button("🤖 Generate Summary Now", key=f"gen_sum_{idx_item}", type="primary"):
                            with st.spinner("Generating summary…"):
                                try:
                                    doc_type_v = (
                                        "PDF document"          if fext_v == "pdf"  else
                                        "PowerPoint presentation" if fext_v == "pptx" else
                                        "document"
                                    )
                                    s = summarise_with_llm(get_llm(), raw_text, doc_type_v)
                                    st.session_state.ingest_log[idx_item]["summary"] = s
                                    st.rerun()
                                except Exception as e:
                                    st.error(f"❌ {e}")

                # ── PDF inline viewer ──
                if fext_v == "pdf" and fb:
                    st.markdown('<div class="fv-section-head">PDF Viewer</div>', unsafe_allow_html=True)
                    b64 = base64.b64encode(fb).decode("utf-8")
                    pdf_html = f"""
                    <!DOCTYPE html>
                    <html>
                    <head><style>
                      body {{ margin: 0; padding: 0; background: #fff; }}
                      iframe {{
                        width: 100%; height: 720px;
                        border: none; display: block;
                        border-radius: 10px;
                      }}
                      .fallback {{
                        padding: 20px; text-align: center;
                        font-family: sans-serif; color: #666;
                      }}
                    </style></head>
                    <body>
                      <iframe
                        src="data:application/pdf;base64,{b64}#toolbar=1&navpanes=1&scrollbar=1"
                        type="application/pdf"
                      >
                        <div class="fallback">
                          <p>Your browser doesn't support inline PDF viewing.</p>
                          <a href="data:application/pdf;base64,{b64}"
                             download="{name}"
                             style="color:#7c5cfc; font-weight:600;">
                            ⬇️ Download PDF instead
                          </a>
                        </div>
                      </iframe>
                    </body>
                    </html>
                    """
                    components.html(pdf_html, height=740, scrolling=False)

                    # Download button below viewer
                    st.download_button(
                        label="⬇️  Download PDF",
                        data=fb,
                        file_name=name,
                        mime="application/pdf",
                        key=f"dl_pdf_{idx_item}",
                        use_container_width=True,
                    )

                # ── PPTX / Text viewer ──
                elif fext_v in ("pptx", "text") and raw_text:
                    st.markdown('<div class="fv-section-head">Extracted Content</div>', unsafe_allow_html=True)
                    preview = raw_text[:6000] + ("\n\n[… truncated — showing first 6 000 chars …]" if len(raw_text) > 6000 else "")
                    st.markdown(f'<div class="fv-text-wrap">{preview}</div>', unsafe_allow_html=True)

                st.markdown("</div>", unsafe_allow_html=True)  # close fv-panel

        st.divider()
        if st.button("🗑️  Clear Ingestion Log", use_container_width=True):
            st.session_state.ingest_log      = []
            st.session_state.ingested_chunks = 0
            st.session_state.kb_view_item    = None
            st.rerun()

    else:
        st.markdown("""
        <div class="il-empty">
          <div class="il-empty-icon">📭</div>
          Nothing ingested yet.<br>Go to <strong>Document RAG</strong> or <strong>YouTube RAG</strong> to add content.
        </div>
        """, unsafe_allow_html=True)


# ─────────────────────────────────────────────
# Footer
# ─────────────────────────────────────────────

st.markdown("""
<div class="sb-footer">
  <b>Second Brain</b> &nbsp;·&nbsp; Streamlit &nbsp;·&nbsp; ChromaDB &nbsp;·&nbsp; LangChain &nbsp;·&nbsp; Groq &nbsp;·&nbsp; pypdf &nbsp;·&nbsp; python-pptx
</div>
""", unsafe_allow_html=True)